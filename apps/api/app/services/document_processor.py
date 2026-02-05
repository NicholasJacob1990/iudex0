"""
Processador de Documentos com Suporte a Contexto Ilimitado
Sistema que divide documentos grandes em chunks e processa em paralelo
"""

from typing import List, Dict, Any, Optional
from dataclasses import dataclass
import hashlib
from loguru import logger

from app.core.config import settings


@dataclass
class DocumentChunk:
    """Chunk de documento"""
    id: str
    content: str
    position: int
    metadata: Dict[str, Any]
    embedding: Optional[List[float]] = None
    
    @property
    def token_count(self) -> int:
        """Estima número de tokens (1 token ≈ 4 caracteres)"""
        return len(self.content) // 4


@dataclass
class ProcessedDocument:
    """Documento processado"""
    document_id: str
    original_size: int
    chunks: List[DocumentChunk]
    total_tokens: int
    extracted_text: str
    metadata: Dict[str, Any]


class DocumentChunker:
    """
    Divide documentos grandes em chunks semânticos
    Mantém contexto entre chunks para não perder informação
    """
    
    def __init__(
        self,
        chunk_size: int = 4000,  # tokens por chunk
        overlap: int = 200,  # overlap entre chunks
    ):
        self.chunk_size = chunk_size
        self.overlap = overlap
        logger.info(f"DocumentChunker inicializado - Chunk size: {chunk_size}, Overlap: {overlap}")
    
    def chunk_by_tokens(self, text: str, metadata: Dict[str, Any] = None) -> List[DocumentChunk]:
        """
        Divide texto em chunks baseado em tokens
        Mantém overlap para preservar contexto
        """
        if metadata is None:
            metadata = {}
        
        # Estimar tokens (1 token ≈ 4 caracteres)
        chars_per_chunk = self.chunk_size * 4
        chars_overlap = self.overlap * 4
        
        chunks: List[DocumentChunk] = []
        position = 0
        start = 0
        
        while start < len(text):
            # Calcular fim do chunk
            end = start + chars_per_chunk
            
            # Se não é o último chunk, tentar quebrar em parágrafo ou frase
            if end < len(text):
                # Procurar quebra de parágrafo
                paragraph_break = text.rfind('\n\n', start + chars_per_chunk - 500, end + 500)
                if paragraph_break > start:
                    end = paragraph_break
                else:
                    # Procurar quebra de frase
                    sentence_break = text.rfind('. ', start + chars_per_chunk - 200, end + 200)
                    if sentence_break > start:
                        end = sentence_break + 1
            
            # Extrair chunk
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunk_id = self._generate_chunk_id(chunk_text, position)
                
                chunks.append(DocumentChunk(
                    id=chunk_id,
                    content=chunk_text,
                    position=position,
                    metadata={
                        **metadata,
                        "start_char": start,
                        "end_char": end,
                        "chunk_index": position,
                    }
                ))
                
                position += 1
            
            # Próximo chunk com overlap
            start = end - chars_overlap if end < len(text) else end
            
            # Evitar loop infinito
            if start >= len(text):
                break
        
        logger.info(f"Documento dividido em {len(chunks)} chunks")
        return chunks
    
    def chunk_by_pages(
        self,
        pages: List[str],
        pages_per_chunk: int = 10,
        metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Divide documento em chunks baseado em páginas
        Útil para PDFs onde queremos manter páginas juntas
        """
        if metadata is None:
            metadata = {}
        
        chunks: List[DocumentChunk] = []
        
        for i in range(0, len(pages), pages_per_chunk):
            chunk_pages = pages[i:i + pages_per_chunk]
            chunk_text = "\n\n".join(chunk_pages)
            
            chunk_id = self._generate_chunk_id(chunk_text, i // pages_per_chunk)
            
            chunks.append(DocumentChunk(
                id=chunk_id,
                content=chunk_text,
                position=i // pages_per_chunk,
                metadata={
                    **metadata,
                    "page_start": i,
                    "page_end": min(i + pages_per_chunk, len(pages)),
                    "page_number": i + 1,
                    "num_pages": len(chunk_pages),
                }
            ))
        
        logger.info(f"Documento dividido em {len(chunks)} chunks por página")
        return chunks
    
    def chunk_semantically(
        self,
        text: str,
        section_markers: List[str] = None,
        metadata: Dict[str, Any] = None
    ) -> List[DocumentChunk]:
        """
        Divide documento em chunks semânticos (por seções)
        Identifica seções baseado em marcadores
        """
        if section_markers is None:
            # Marcadores padrão para documentos jurídicos
            section_markers = [
                "RELATÓRIO",
                "FUNDAMENTAÇÃO",
                "VOTO",
                "DECISÃO",
                "DISPOSITIVO",
                "EMENTA",
                "DOS FATOS",
                "DO DIREITO",
                "I -",
                "II -",
                "III -",
            ]
        
        if metadata is None:
            metadata = {}
        
        chunks: List[DocumentChunk] = []
        current_section = ""
        current_content = []
        position = 0
        
        lines = text.split('\n')
        
        for line in lines:
            # Verificar se linha é um marcador de seção
            is_section = any(line.strip().upper().startswith(marker) for marker in section_markers)
            
            if is_section and current_content:
                # Salvar seção anterior
                section_text = '\n'.join(current_content).strip()
                if section_text:
                    chunk_id = self._generate_chunk_id(section_text, position)
                    chunks.append(DocumentChunk(
                        id=chunk_id,
                        content=section_text,
                        position=position,
                        metadata={
                            **metadata,
                            "section": current_section,
                            "semantic_chunk": True,
                        }
                    ))
                    position += 1
                
                # Iniciar nova seção
                current_section = line.strip()
                current_content = [line]
            else:
                current_content.append(line)
        
        # Adicionar última seção
        if current_content:
            section_text = '\n'.join(current_content).strip()
            if section_text:
                chunk_id = self._generate_chunk_id(section_text, position)
                chunks.append(DocumentChunk(
                    id=chunk_id,
                    content=section_text,
                    position=position,
                    metadata={
                        **metadata,
                        "section": current_section,
                        "semantic_chunk": True,
                    }
                ))
        
        # Se ficou muito grande, dividir novamente por tokens
        final_chunks = []
        for chunk in chunks:
            if chunk.token_count > self.chunk_size:
                sub_chunks = self.chunk_by_tokens(
                    chunk.content,
                    {**chunk.metadata, "parent_chunk_id": chunk.id}
                )
                final_chunks.extend(sub_chunks)
            else:
                final_chunks.append(chunk)
        
        logger.info(f"Documento dividido em {len(final_chunks)} chunks semânticos")
        return final_chunks
    
    def _generate_chunk_id(self, content: str, position: int) -> str:
        """Gera ID único para chunk"""
        content_hash = hashlib.md5(content.encode()).hexdigest()[:8]
        return f"chunk-{position}-{content_hash}"


class UnlimitedContextProcessor:
    """
    Processador que suporta contexto ilimitado
    
    Estratégias:
    1. Map-Reduce: Processa chunks em paralelo e consolida
    2. Hierarchical: Cria resumos hierárquicos
    3. Rolling: Mantém janela deslizante de contexto
    """
    
    def __init__(self):
        self.chunker = DocumentChunker()
        logger.info("UnlimitedContextProcessor inicializado")
    
    async def process_large_document(
        self,
        text: str,
        task: str,
        strategy: str = "map-reduce",
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Processa documento grande com estratégia escolhida
        
        Args:
            text: Texto completo do documento
            task: Tarefa a realizar (resumir, analisar, etc)
            strategy: "map-reduce", "hierarchical", ou "rolling"
            metadata: Metadados adicionais
        """
        logger.info(f"Processando documento grande - Estratégia: {strategy}, Tamanho: {len(text)} chars")
        
        if strategy == "map-reduce":
            return await self._process_map_reduce(text, task, metadata)
        elif strategy == "hierarchical":
            return await self._process_hierarchical(text, task, metadata)
        elif strategy == "rolling":
            return await self._process_rolling(text, task, metadata)
        else:
            raise ValueError(f"Estratégia desconhecida: {strategy}")
    
    async def _process_map_reduce(
        self,
        text: str,
        task: str,
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Map-Reduce: Processa cada chunk independentemente e depois consolida
        Ideal para: resumos, extração de informações
        """
        logger.info("Executando estratégia Map-Reduce")
        
        # 1. Dividir em chunks
        chunks = self.chunker.chunk_by_tokens(text, metadata)
        
        # 2. MAP: Processar cada chunk (em paralelo)
        # TODO: Implementar processamento paralelo com asyncio.gather
        chunk_results = []
        for chunk in chunks:
            # Aqui chamaria o agente IA para processar o chunk
            # result = await ai_agent.process(chunk.content, task)
            result = {
                "chunk_id": chunk.id,
                "position": chunk.position,
                "summary": f"[Resumo do chunk {chunk.position}]",  # Placeholder
            }
            chunk_results.append(result)
        
        # 3. REDUCE: Consolidar resultados
        # TODO: Chamar IA para consolidar todos os resumos
        final_result = {
            "strategy": "map-reduce",
            "total_chunks": len(chunks),
            "chunk_results": chunk_results,
            "consolidated": "[Resultado consolidado de todos os chunks]",  # Placeholder
            "metadata": {
                "original_size": len(text),
                "total_tokens": sum(c.token_count for c in chunks),
            }
        }
        
        logger.info(f"Map-Reduce concluído - {len(chunks)} chunks processados")
        return final_result
    
    async def _process_hierarchical(
        self,
        text: str,
        task: str,
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Hierarchical: Cria resumos em níveis, cada nível resume o anterior
        Ideal para: análise detalhada, manutenção de contexto global
        """
        logger.info("Executando estratégia Hierarchical")
        
        levels = []
        current_text = text
        level = 0
        
        # Continuar até conseguir processar em um único chunk
        while len(current_text) > self.chunker.chunk_size * 4:
            logger.info(f"Nível {level}: {len(current_text)} chars")
            
            # Dividir em chunks
            chunks = self.chunker.chunk_by_tokens(current_text, metadata)
            
            # Resumir cada chunk
            summaries = []
            for chunk in chunks:
                # TODO: Chamar IA para resumir
                summary = f"[Resumo do chunk {chunk.position}]"  # Placeholder
                summaries.append(summary)
            
            # Consolidar resumos deste nível
            current_text = "\n\n".join(summaries)
            levels.append({
                "level": level,
                "num_chunks": len(chunks),
                "summaries": summaries,
            })
            
            level += 1
            
            # Proteção contra loop infinito
            if level > 10:
                break
        
        # Processar nível final
        # TODO: Chamar IA para análise final
        final_analysis = "[Análise final do documento completo]"  # Placeholder
        
        result = {
            "strategy": "hierarchical",
            "num_levels": len(levels),
            "levels": levels,
            "final_analysis": final_analysis,
            "metadata": {
                "original_size": len(text),
            }
        }
        
        logger.info(f"Hierarchical concluído - {len(levels)} níveis processados")
        return result
    
    async def _process_rolling(
        self,
        text: str,
        task: str,
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Rolling Window: Processa com janela deslizante mantendo contexto
        Ideal para: geração de documento contínuo, manutenção de narrativa
        """
        logger.info("Executando estratégia Rolling Window")
        
        chunks = self.chunker.chunk_by_tokens(text, metadata)
        
        results = []
        context_window = []
        window_size = 3  # Manter últimos 3 chunks no contexto
        
        for i, chunk in enumerate(chunks):
            logger.info(f"Processando chunk {i+1}/{len(chunks)} com contexto")
            
            # Montar contexto (chunks anteriores)
            context = "\n\n".join([c.content for c in context_window])
            
            # Processar chunk com contexto
            # TODO: Chamar IA com contexto
            result = {
                "chunk_id": chunk.id,
                "position": chunk.position,
                "context_size": len(context_window),
                "result": f"[Resultado do chunk {i} com contexto]",  # Placeholder
            }
            results.append(result)
            
            # Atualizar janela de contexto
            context_window.append(chunk)
            if len(context_window) > window_size:
                context_window.pop(0)
        
        # Consolidar resultados mantendo a narrativa
        final_result = {
            "strategy": "rolling",
            "total_chunks": len(chunks),
            "window_size": window_size,
            "results": results,
            "consolidated": "[Resultado consolidado mantendo narrativa]",  # Placeholder
            "metadata": {
                "original_size": len(text),
                "total_tokens": sum(c.token_count for c in chunks),
            }
        }
        
        logger.info(f"Rolling Window concluído - {len(chunks)} chunks processados")
        return final_result


import pdfplumber
from docx import Document as DocxDocument
import pytesseract
from PIL import Image
import os

# Funções auxiliares para extração de texto de diferentes formatos

@dataclass
class PageText:
    """Texto extraído com metadados de página."""
    page_number: int
    text: str
    line_start: int = 0
    line_end: int = 0


async def extract_text_from_pdf(file_path: str) -> str:
    """
    Extrai texto de PDF usando pdfplumber
    Melhor para manter layout e extrair tabelas
    """
    logger.info(f"Extraindo texto de PDF: {file_path}")
    text_content = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                text_content.append(text)

        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Erro ao extrair texto do PDF {file_path}: {e}")
        # Fallback para pypdf se necessário ou re-raise
        raise


async def extract_pages_from_pdf(file_path: str) -> List[PageText]:
    """
    Extrai texto de PDF com metadados de página e linha.
    Retorna lista de PageText com page_number, line_start e line_end.
    """
    logger.info(f"Extraindo páginas de PDF com metadados: {file_path}")
    pages: List[PageText] = []
    global_line = 0

    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                line_count = text.count("\n") + 1 if text.strip() else 0
                pages.append(PageText(
                    page_number=page.page_number,
                    text=text,
                    line_start=global_line,
                    line_end=global_line + line_count - 1 if line_count > 0 else global_line,
                ))
                global_line += line_count

        logger.info(f"PDF extraído: {len(pages)} páginas, {global_line} linhas totais")
        return pages
    except Exception as e:
        logger.error(f"Erro ao extrair páginas do PDF {file_path}: {e}")
        raise


async def extract_paragraphs_from_docx(file_path: str) -> List[PageText]:
    """
    Extrai parágrafos de DOCX com índice de parágrafo como referência de linha.
    Retorna lista de PageText (page_number = índice do parágrafo).
    """
    logger.info(f"Extraindo parágrafos de DOCX com metadados: {file_path}")
    paragraphs: List[PageText] = []

    try:
        doc = DocxDocument(file_path)
        line_offset = 0

        for idx, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue
            line_count = text.count("\n") + 1
            paragraphs.append(PageText(
                page_number=idx + 1,  # índice do parágrafo como referência
                text=text,
                line_start=line_offset,
                line_end=line_offset + line_count - 1,
            ))
            line_offset += line_count

        # Tabelas
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    line_offset += 1
                    paragraphs.append(PageText(
                        page_number=len(doc.paragraphs) + 1,
                        text=row_text,
                        line_start=line_offset,
                        line_end=line_offset,
                    ))

        logger.info(f"DOCX extraído: {len(paragraphs)} blocos de texto")
        return paragraphs
    except Exception as e:
        logger.error(f"Erro ao extrair parágrafos do DOCX {file_path}: {e}")
        raise

async def extract_text_from_docx(file_path: str) -> str:
    """Extrai texto de DOCX usando python-docx"""
    logger.info(f"Extraindo texto de DOCX: {file_path}")
    try:
        doc = DocxDocument(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
                
        # Também extrair texto de tabelas
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text_content.append(" | ".join(row_text))
                
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Erro ao extrair texto do DOCX {file_path}: {e}")
        raise

async def extract_text_from_image(file_path: str, use_hybrid: bool = True) -> str:
    """
    Extrai texto de imagem usando OCR

    Args:
        file_path: Caminho da imagem
        use_hybrid: Se True, usa serviço híbrido com fallback para cloud
    """
    logger.info(f"Extraindo texto de imagem: {file_path}")

    if use_hybrid:
        try:
            from app.services.ocr_service import get_ocr_service

            ocr_service = get_ocr_service()
            result = await ocr_service.extract_text_from_image(file_path)

            if result.error:
                logger.warning(f"OCR híbrido falhou, usando Tesseract direto: {result.error}")
                # Fallback para Tesseract direto
                image = Image.open(file_path)
                return pytesseract.image_to_string(image, lang='por')

            logger.info(f"OCR concluído via {result.provider.value}")
            return result.text
        except ImportError:
            logger.warning("Serviço OCR híbrido não disponível, usando Tesseract")

    # Fallback: Tesseract direto
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='por')
        return text
    except Exception as e:
        logger.error(f"Erro ao extrair texto da imagem {file_path}: {e}")
        return f"[Erro no OCR: {str(e)}]"


async def extract_text_from_pdf_with_ocr(
    file_path: str,
    use_hybrid: bool = True,
    force_ocr: bool = False,
) -> str:
    """
    Extrai texto de PDF usando estratégia híbrida de OCR

    Estratégia:
    1. Se PDF tem texto selecionável e force_ocr=False → pdfplumber
    2. Se volume baixo → Tesseract local (gratuito)
    3. Se volume alto ou Tesseract falha → Cloud OCR (Azure/Google/Gemini)

    Args:
        file_path: Caminho do PDF
        use_hybrid: Se True, usa serviço híbrido inteligente
        force_ocr: Se True, força OCR mesmo em PDFs com texto selecionável
    """
    logger.info(f"Extraindo texto de PDF: {file_path} (hybrid={use_hybrid}, force_ocr={force_ocr})")

    if use_hybrid:
        try:
            from app.services.ocr_service import get_ocr_service

            ocr_service = get_ocr_service()
            result = await ocr_service.extract_text_from_pdf(
                file_path,
                force_ocr=force_ocr,
            )

            if result.error:
                logger.warning(f"OCR híbrido falhou: {result.error}")
                # Fallback para implementação original
                return await _extract_text_from_pdf_tesseract(file_path)

            logger.info(
                f"Extração concluída via {result.provider.value}: "
                f"{result.pages_processed} páginas, {len(result.text)} chars"
            )
            return result.text
        except ImportError as e:
            logger.warning(f"Serviço OCR híbrido não disponível: {e}")

    # Fallback: implementação original com Tesseract
    return await _extract_text_from_pdf_tesseract(file_path)


async def _extract_text_from_pdf_tesseract(file_path: str) -> str:
    """
    Implementação original de OCR com Tesseract (fallback)
    """
    logger.info(f"Aplicando OCR Tesseract em PDF: {file_path}")
    try:
        from pdf2image import convert_from_path

        logger.info("Convertendo PDF para imagens...")
        images = convert_from_path(file_path, dpi=300)

        logger.info(f"PDF convertido em {len(images)} imagens")

        ocr_texts = []
        for i, image in enumerate(images, 1):
            logger.info(f"Aplicando OCR na página {i}/{len(images)}")
            page_text = pytesseract.image_to_string(image, lang='por')

            if page_text.strip():
                ocr_texts.append(f"--- Página {i} ---\n{page_text}")
            else:
                ocr_texts.append(f"--- Página {i} ---\n[Página sem texto detectado]")

        result = "\n\n".join(ocr_texts)
        logger.info(f"OCR concluído: {len(result)} caracteres extraídos")
        return result

    except ImportError as e:
        logger.error("Bibliotecas necessárias não instaladas")
        logger.error("Instale com: pip install pdf2image")
        logger.error("macOS: brew install poppler")
        logger.error("Linux: apt-get install poppler-utils")
        return f"[Erro: bibliotecas de OCR não instaladas - {str(e)}]"
    except Exception as e:
        logger.error(f"Erro ao aplicar OCR em PDF {file_path}: {e}")
        return f"[Erro no OCR do PDF: {str(e)}]"



async def extract_text_from_odt(file_path: str) -> str:
    """Extrai texto de ODT (OpenDocument Text)"""
    logger.info(f"Extraindo texto de ODT: {file_path}")
    try:
        from odf import text, teletype
        from odf.opendocument import load
        
        textdoc = load(file_path)
        allparas = textdoc.getElementsByType(text.P)
        text_content = []
        
        for para in allparas:
            para_text = teletype.extractText(para)
            if para_text.strip():
                text_content.append(para_text)
        
        return "\n\n".join(text_content)
    except ImportError:
        logger.error("Biblioteca odfpy não instalada. Instale com: pip install odfpy")
        return "[Erro: biblioteca odfpy não instalada]"
    except Exception as e:
        logger.error(f"Erro ao extrair texto do ODT {file_path}: {e}")
        return f"[Erro na extração ODT: {str(e)}]"


async def extract_text_from_pptx(file_path: str) -> str:
    """
    Extrai texto de PPTX (PowerPoint) usando python-pptx
    Extrai texto de slides, notas e tabelas, mantendo ordem dos slides
    """
    logger.info(f"Extraindo texto de PPTX: {file_path}")
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_content: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: List[str] = []

            # Extrair texto de todas as shapes com text_frame
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_texts.append(para_text)

                # Extrair texto de tabelas
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            slide_texts.append(row_text)

            # Extrair notas do slide
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[Notas] {notes_text}")

            # Pular slides vazios
            if slide_texts:
                text_content.append(
                    f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts)
                )

        result = "\n\n".join(text_content)
        logger.info(f"PPTX extraído: {len(prs.slides)} slides, {len(result)} chars")
        return result
    except ImportError:
        logger.error("Biblioteca python-pptx não instalada. Instale com: pip install python-pptx")
        return "[Erro: biblioteca python-pptx não instalada]"
    except Exception as e:
        logger.error(f"Erro ao extrair texto do PPTX {file_path}: {e}")
        return f"[Erro na extração PPTX: {str(e)}]"


async def extract_text_from_xlsx(file_path: str) -> str:
    """
    Extrai texto de XLSX (Excel) usando openpyxl
    Formata com cabeçalhos de planilha e separadores de célula
    """
    logger.info(f"Extraindo texto de XLSX: {file_path}")
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_content: List[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_lines: List[str] = []

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                # Converter células para string, tratar None como vazio
                cell_values = [
                    str(cell) if cell is not None else ""
                    for cell in row
                ]
                # Pular linhas completamente vazias
                if not any(v.strip() for v in cell_values):
                    continue
                row_text = " | ".join(cell_values)
                sheet_lines.append(row_text)

            if sheet_lines:
                text_content.append(
                    f"=== Sheet: {sheet_name} ===\n" + "\n".join(sheet_lines)
                )

        num_sheets = len(wb.sheetnames)
        wb.close()
        result = "\n\n".join(text_content)
        logger.info(f"XLSX extraído: {num_sheets} planilhas, {len(result)} chars")
        return result
    except ImportError:
        logger.error("Biblioteca openpyxl não instalada. Instale com: pip install openpyxl")
        return "[Erro: biblioteca openpyxl não instalada]"
    except Exception as e:
        logger.error(f"Erro ao extrair texto do XLSX {file_path}: {e}")
        return f"[Erro na extração XLSX: {str(e)}]"


async def extract_text_from_csv(file_path: str) -> str:
    """
    Extrai texto de CSV usando módulo csv padrão
    Auto-detecta delimitador com csv.Sniffer
    """
    logger.info(f"Extraindo texto de CSV: {file_path}")
    try:
        import csv

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            # Ler amostra para detectar delimitador
            sample = f.read(8192)
            f.seek(0)

            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                # Fallback para delimitador padrão (vírgula)
                dialect = csv.excel

            reader = csv.reader(f, dialect)
            lines: List[str] = []

            for row in reader:
                cell_values = [cell.strip() for cell in row]
                if any(cell_values):
                    lines.append(" | ".join(cell_values))

        result = "\n".join(lines)
        logger.info(f"CSV extraído: {len(lines)} linhas, {len(result)} chars")
        return result
    except Exception as e:
        logger.error(f"Erro ao extrair texto do CSV {file_path}: {e}")
        return f"[Erro na extração CSV: {str(e)}]"


async def extract_text_from_rtf(file_path: str) -> str:
    """
    Extrai texto de RTF usando regex para strip de control words
    Não requer dependências externas
    """
    logger.info(f"Extraindo texto de RTF: {file_path}")
    try:
        import re

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # Remover grupos de controle {\*...}
        text = re.sub(r'\{\\\*[^}]*\}', '', content)
        # Remover control words com argumentos (ex: \fs24, \par)
        text = re.sub(r'\\[a-zA-Z]+\d*\s?', ' ', text)
        # Remover caracteres de escape RTF (ex: \', \\)
        text = re.sub(r'\\[^a-zA-Z]', '', text)
        # Remover chaves restantes
        text = re.sub(r'[{}]', '', text)
        # Limpar espaços múltiplos
        text = re.sub(r'[ \t]+', ' ', text)
        # Limpar linhas em branco múltiplas
        text = re.sub(r'\n\s*\n', '\n\n', text)

        result = text.strip()
        logger.info(f"RTF extraído: {len(result)} chars")
        return result
    except Exception as e:
        logger.error(f"Erro ao extrair texto do RTF {file_path}: {e}")
        return f"[Erro na extração RTF: {str(e)}]"


async def extract_text_from_zip(file_path: str) -> dict:
    """
    Extrai e processa arquivos de um ZIP
    Retorna dicionário com informações dos arquivos extraídos
    """
    logger.info(f"Processando arquivo ZIP: {file_path}")
    try:
        import zipfile
        import tempfile
        
        results = {
            "files": [],
            "total_files": 0,
            "extracted_text": "",
            "errors": []
        }
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            # Criar diretório temporário para extração
            with tempfile.TemporaryDirectory() as temp_dir:
                zip_ref.extractall(temp_dir)
                
                # Listar arquivos extraídos
                file_list = zip_ref.namelist()
                results["total_files"] = len(file_list)
                
                all_text = []
                
                # Processar cada arquivo
                for file_name in file_list:
                    extracted_path = os.path.join(temp_dir, file_name)
                    
                    # Pular diretórios
                    if os.path.isdir(extracted_path):
                        continue
                    
                    file_info = {
                        "name": file_name,
                        "size": os.path.getsize(extracted_path),
                        "status": "processed"
                    }
                    
                    # Tentar extrair texto baseado na extensão
                    ext = os.path.splitext(file_name)[1].lower()
                    
                    try:
                        if ext == '.pdf':
                            text = await extract_text_from_pdf(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext in ['.docx', '.doc']:
                            text = await extract_text_from_docx(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext == '.odt':
                            text = await extract_text_from_odt(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext == '.pptx':
                            text = await extract_text_from_pptx(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext in ['.xlsx', '.xls']:
                            text = await extract_text_from_xlsx(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext == '.csv':
                            text = await extract_text_from_csv(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext == '.rtf':
                            text = await extract_text_from_rtf(extracted_path)
                            all_text.append(f"=== {file_name} ===\n{text}")
                        elif ext == '.txt':
                            with open(extracted_path, 'r', encoding='utf-8', errors='ignore') as f:
                                text = f.read()
                                all_text.append(f"=== {file_name} ===\n{text}")
                        else:
                            file_info["status"] = "unsupported"
                    except Exception as e:
                        logger.error(f"Erro ao processar {file_name}: {e}")
                        file_info["status"] = "error"
                        file_info["error"] = str(e)
                        results["errors"].append(f"{file_name}: {str(e)}")
                    
                    results["files"].append(file_info)
                
                results["extracted_text"] = "\n\n".join(all_text)
        
        logger.info(f"ZIP processado: {len(results['files'])} arquivos")
        return results
        
    except Exception as e:
        logger.error(f"Erro ao processar ZIP {file_path}: {e}")
        return {
            "error": str(e),
            "files": [],
            "total_files": 0,
            "extracted_text": ""
        }


async def transcribe_audio_video(file_path: str, media_type: str = "audio") -> str:
    """
    Transcreve áudio ou vídeo usando Whisper (OpenAI)
    
    Args:
        file_path: Caminho do arquivo de áudio/vídeo
        media_type: "audio" ou "video"
    """
    logger.info(f"Transcrevendo {media_type}: {file_path}")
    try:
        # Verificar se OpenAI Whisper está disponível
        import openai
        
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            logger.warning("OPENAI_API_KEY não configurada, tentando whisper local")
            return await _transcribe_with_local_whisper(file_path)
        
        # Usar API do OpenAI Whisper
        openai.api_key = api_key
        
        with open(file_path, "rb") as audio_file:
            transcript = openai.Audio.transcribe(
                model="whisper-1",
                file=audio_file,
                language="pt"  # Português
            )
        
        return transcript.text
        
    except ImportError:
        logger.warning("openai não instalado, tentando whisper local")
        return await _transcribe_with_local_whisper(file_path)
    except Exception as e:
        logger.error(f"Erro na transcrição: {e}")
        return f"[Erro na transcrição: {str(e)}]"


async def _transcribe_with_local_whisper(file_path: str) -> str:
    """
    Transcreve usando Whisper local (faster-whisper ou whisper)
    """
    try:
        import whisper
        
        logger.info("Usando Whisper local para transcrição")
        model = whisper.load_model("base")  # ou "small", "medium", "large"
        result = model.transcribe(file_path, language="pt")
        
        return result["text"]
        
    except ImportError:
        logger.error("whisper não instalado. Instale com: pip install openai-whisper")
        return "[Erro: biblioteca whisper não instalada]"
    except Exception as e:
        logger.error(f"Erro no Whisper local: {e}")
        return f"[Erro na transcrição local: {str(e)}]"
